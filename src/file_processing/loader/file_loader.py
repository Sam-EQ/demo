import os
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
import ezdxf
import subprocess
import os
from docx import Document

# def load_pdf(path):
#     reader = PdfReader(path)
#     return "\n".join([page.extract_text() or "" for page in reader.pages])


def extract_with_pdfplumber(pdf_path, output_dir="pdfplumber_images"):
    import os
    os.makedirs(output_dir, exist_ok=True)
    
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            # Extract images
            for j, img in enumerate(page.images):
                # pdfplumber returns image metadata
                print(f"Image {j} on page {i+1}: {img}")
            
            # Or render the whole page
            page_img = page.to_image(resolution=300)
            page_img.save(f"{output_dir}/page_{i+1}.png")


def load_doc(path: str) -> str:
    subprocess.run(
        [   
            "libreoffice",
            "--headless",
            "--convert-to",
            "docx",
            path,
            "--outdir",
            os.path.dirname(path)
        ],
        check=True
    )
    docx_path = path.replace(".doc", ".docx")
    doc = Document(docx_path)
    return "\n".join(p.text for p in doc.paragraphs)

def load_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def load_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def load_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)


def load_excel(path):
    df = pd.read_excel(path)
    return df.to_string(index=False)


def load_odf(path):
    doc = load(path)
    return teletype.extractText(doc)


def load_dwg(path):
    doc = ezdxf.readfile(path)
    msp = doc.modelspace()
    texts = [e.dxf.text for e in msp.query("TEXT")]
    return "\n".join(texts)

def load_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    loaders = {
        ".pdf": load_pdf,
        ".docx": load_docx,
        ".doc":load_doc,
        ".pptx": load_pptx,
        ".csv": load_csv,
        ".xls": load_excel,
        ".xlsx": load_excel,
        ".odf": load_odf,
        ".dwg": load_dwg,
    }

    if ext not in loaders:
        raise ValueError(f"Unsupported file type: {ext}")

    return loaders[ext](file_path)


if __name__ == "__main__":
    print("===== testing file loader =====")
    value = extract_with_pdfplumber("/Users/ezhilrajselvaraj/Downloads/test.pdf")
    # print(value["text"])
    with open("ok.txt",'w') as f:
        f.write(value['text'])
    # print(value['images'])
